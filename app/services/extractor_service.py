from __future__ import annotations

import re
from pathlib import Path
from typing import Iterable

from pypdf import PdfReader
from pptx import Presentation
from pptx.shapes.base import BaseShape

from app.core.exceptions import EmptyDocumentError, UnsupportedFileTypeError
from app.models.domain import PageText


class ExtractorService:
    def extract(self, file_path: Path) -> list[PageText]:
        extension = file_path.suffix.lower()
        if extension == ".pdf":
            pages = self._extract_pdf(file_path)
        elif extension == ".pptx":
            pages = self._extract_pptx(file_path)
        else:
            raise UnsupportedFileTypeError(f"Không hỗ trợ định dạng {extension}")

        pages = [page for page in pages if page.text.strip()]
        if not pages:
            raise EmptyDocumentError(
                "Không trích xuất được văn bản. Slide dạng ảnh cần bổ sung OCR."
            )
        return pages

    def _extract_pdf(self, file_path: Path) -> list[PageText]:
        reader = PdfReader(str(file_path))
        return [
            PageText(
                page_number=index,
                text=self._normalize_text(page.extract_text() or ""),
            )
            for index, page in enumerate(reader.pages, start=1)
        ]

    def _extract_pptx(self, file_path: Path) -> list[PageText]:
        presentation = Presentation(str(file_path))
        pages: list[PageText] = []

        for page_number, slide in enumerate(presentation.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                parts.extend(self._extract_shape_text(shape))
            pages.append(
                PageText(
                    page_number=page_number,
                    text=self._normalize_text("\n".join(parts)),
                )
            )
        return pages

    def _extract_shape_text(self, shape: BaseShape) -> Iterable[str]:
        if getattr(shape, "has_text_frame", False):
            text = getattr(shape, "text", "").strip()
            if text:
                yield text

        if getattr(shape, "has_table", False):
            table = shape.table
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if values:
                    yield " | ".join(values)

        if getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
            for child in shape.shapes:
                yield from self._extract_shape_text(child)

    @staticmethod
    def _normalize_text(text: str) -> str:
        text = text.replace("\x00", " ")
        lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines()]
        return "\n".join(line for line in lines if line)
